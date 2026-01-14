# core/builder.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from config import Config

class PPTXBuilder:
    def __init__(self):
        self.prs = Presentation()
        # Remove default blank slide
        if len(self.prs.slides) > 0:
            xml_slides = self.prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[0])

    def _pixels_to_inches(self, px_val, img_dim_px, slide_dim_inch):
        """
        Converts pixel coordinate to PPTX inches using ratio mapping.
        Formula: (Value_Px / Total_Img_Px) * Total_Slide_Inches
        """
        return (px_val / img_dim_px) * slide_dim_inch

    def _pixels_to_points(self, px_val, img_h_px, slide_h_inch):
        """
        Converts pixel font size to PPTX points.
        1 inch = 72 points.
        """
        return (px_val / img_h_px) * slide_h_inch * 72

    def _sanitize_text(self, text):
        """
        Removes control characters that break PPTX XML.
        Valid XML characters are: #x9 | #xA | #xD | [#x20-#xD7FF] | [#xE000-#xFFFD] | [#x10000-#x10FFFF]
        """
        if not text:
            return ""
        # Keep only valid characters
        return "".join(ch for ch in text if ord(ch) in [0x09, 0x0A, 0x0D] or (0x20 <= ord(ch) <= 0xD7FF) or (0xE000 <= ord(ch) <= 0xFFFD))

    def create_slide(self, original_image: Image.Image, elements: list, cleaned_image: Image.Image, page_num: int):
        """
        Creates a slide with three layers: Original -> Cleaned -> Text Overlay.
        """
        img_w, img_h = original_image.size
        
        # 1. Configure Slide Dimensions
        slide_width_inch = 10.0
        aspect_ratio = img_h / img_w
        slide_height_inch = slide_width_inch * aspect_ratio
        
        self.prs.slide_width = Inches(slide_width_inch)
        self.prs.slide_height = Inches(slide_height_inch)
        
        # 2. Add a blank slide layout
        blank_slide_layout = self.prs.slide_layouts[6] 
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # 3. Layer 0: Original Unmasked Image (Bottom)
        orig_path = os.path.join(Config.TEMP_DIR, f"orig_p{page_num}.png")
        original_image.save(orig_path)
        slide.shapes.add_picture(orig_path, 0, 0, Inches(slide_width_inch), Inches(slide_height_inch))

        # 4. Layer 1: Cleaned Background image (Middle)
        bg_path = os.path.join(Config.TEMP_DIR, f"bg_p{page_num}.png")
        cleaned_image.save(bg_path)
        slide.shapes.add_picture(bg_path, 0, 0, Inches(slide_width_inch), Inches(slide_height_inch))

        # 5. Layer 2+: Process Elements (Top Overlay)
        for elem in elements:
            if elem['type'] not in ['text', 'title', 'header', 'footer']:
                continue

            content = self._sanitize_text(elem.get('content', ''))
            if not content: continue
            
            # Calculate PPTX coordinates
            bbox = elem['bbox']
            x1, y1, x2, y2 = bbox
            w_px = x2 - x1
            h_px = y2 - y1
            
            # Guard against invalid dimensions
            left_val = self._pixels_to_inches(x1, img_w, slide_width_inch)
            top_val = self._pixels_to_inches(y1, img_h, slide_height_inch)
            width_val = self._pixels_to_inches(w_px, img_w, slide_width_inch)
            height_val = self._pixels_to_inches(h_px, img_h, slide_height_inch)
            
            # Constraints: All values must be positive, width/height at least min size
            left = Inches(max(0, left_val))
            top = Inches(max(0, top_val))
            width = Inches(max(0.01, width_val))
            height = Inches(max(0.01, height_val))
            
            # Insert Text Box
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = content
            
            # Precise Font Sizing
            font_px = elem.get('font_size_px', 0)
            if font_px > 0:
                font_pt = self._pixels_to_points(font_px, img_h, slide_height_inch)
                # Clamp font size for sanity and valid PPTX range (1-4000)
                p.font.size = Pt(min(max(font_pt, 1), 400))
            else:
                # Fallback to height-based heuristic if font_px not provided
                approx_pt = (height_val * 72)
                p.font.size = Pt(min(max(approx_pt * 0.8, 8), 40))

    def save(self, output_path: str):
        self.prs.save(output_path)
        print(f"[SUCCESS] Saved PPTX to {output_path}")

